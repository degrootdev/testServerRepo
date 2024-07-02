from glob import glob
import pandas as pd
import os
from os import stat
from pptx import Presentation
from typing import Generator
from datetime import datetime
from os.path import splitext
from spire.doc import Document
from spire.pdf import PdfDocument, PdfTextExtractOptions, PdfTextExtractor
from outlook_msg import Message
from os.path import basename
import numpy as np

import pyarrow as pa
import pyarrow.parquet as pq
import re
import string
import xgboost
from sklearn.feature_extraction.text import HashingVectorizer
from sklearn.metrics import f1_score

class DocumentParser():
    def __init__(self):

        # Setup pdf extraction options
        self.pdf_extractOptions = PdfTextExtractOptions()
        self.pdf_extractOptions.IsSimpleExtraction = True
        

    def load_docx(self, path):
        # Create a Document object
        document = Document()
        # Load a Word document
        document.LoadFromFile(path)
        
        # Extract the text of the document
        document_text = document.GetText()
        
        document_text = document_text.replace("Evaluation Warning: The document was created with Spire.Doc for Python.\r\n","")
        return document_text
    
    def load_pdf(self, path):
        # Create a PdfDocument object
        doc = PdfDocument()
        # Load a PDF document
        doc.LoadFromFile(path)

        # Create an empty list to store extracted text
        test_list = []
        
        # Loop through the pages in the document
        for i in range(doc.Pages.Count):
            page = doc.Pages[i]

            textExtractor = PdfTextExtractor(page)
            text = textExtractor.ExtractText(self.pdf_extractOptions)
            
            text = text.replace("Evaluation Warning : The document was created with Spire.PDF for Python.","").lstrip()
            test_list.append(text)
    
        return "\n".join(test_list) 

    def load_powerpoint(self, filepath):
        with open(filepath, 'rb') as f:
            prs = Presentation(f)
            text_runs = []
        
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
        
        return " ".join(text_runs)

    def load_excel(self, filepath):
        sheets_dict = pd.read_excel(filepath, sheet_name=None)
        sheets_text = []
        for name, sheet in sheets_dict.items():   
            sheet = sheet.replace(np.nan, '', regex=True)
            sheet_text = sheet.to_string(header=True, index=False)
            sheets_text.append(sheet_text)
        return " ".join(sheets_text)

    def load_msg(self, filepath):
        with open(filepath) as msg_file:
            msg = Message(msg_file)
        return msg.body

doc_parser = DocumentParser()


def timeConvert(atime):
    dt = datetime.fromtimestamp(atime)
    return dt.strftime("%d-%m-%Y %H:%M:%S")

def sizeFormat(size):
    newform = format(size/1024, ".2f")
    return newform + " KB"

def get_metadata(filepath, root_folder):
    file_info = stat(filepath)
    
    path_from_root = filepath.split(root_folder)[1].split(os.sep, 1)[-1]

    metadata = {
        'filepath': os.path.join(root_folder, path_from_root),
        'filename': basename(filepath),
        'size': sizeFormat(file_info.st_size),
        'creation_date': timeConvert(file_info.st_ctime),
        'modified_date': timeConvert(file_info.st_mtime),
        # 'last_access_date': timeConvert(file_info.st_atime),   
    }
    return metadata

def get_extension(filepath):
    return splitext(filepath)[-1]

def get_content(filepath, parsers):
    file_extension = get_extension(filepath).lower()
    try:
        # select a parser for the given file extension
        parser = parsers[file_extension]
        
        # parse the file's content
        content = parser(filepath)
        
    except KeyError:
        # file extension doesn't have a parser
        return ''

    return content

def label_generator(target_filename, label_folder, label_mapper, root_folder=None):
    # determine root_folder if it was not already set (else it changes at each recursive function call) # TODO: move out of function
    root_folder = label_folder.split(os.sep)[-1] if not root_folder else root_folder

    with os.scandir(label_folder) as scan:
        for item in scan:
            labeled_filepath = item.path

            # if the item is a file
            if item.is_file():
                labeled_filename = os.path.basename(labeled_filepath)

                # file has the target filename
                if labeled_filename == target_filename:
                    label_str = labeled_filepath.split(root_folder)[1].split(os.sep)[1]
                    yield label_mapper[label_str]

                # # file does not have the target filename
                else:
                    yield '0'
                
            # else if the item is a folder
            else:
                for subitem in label_generator(target_filename, labeled_filepath, label_mapper, root_folder):
                    yield subitem
    

def get_label(filepath, label_folder, label_mapper, batch_size=100):
    # determine the filename of the target file we want to find the label for
    target_filename = os.path.basename(filepath)
    
    # initialize label_generator to help find the correct label(s) in the label_folder
    label_gen = label_generator(target_filename, label_folder, label_mapper)

    all_labels = []
    while True:
        # perform a batch iteration
        batch_labels = list(next(batch_generator(label_gen, batch_size=batch_size)))

        # extend all_labels with the non-duplicate labels found for this file in the current batch
        all_labels.extend(list(set(batch_labels)))

        # whole label_folder is processed
        if not batch_labels:
            # remove duplicate labels
            all_labels = list(set(all_labels))

            # remove '0' (no label) as a label if there were any other labels found for this file
            if len(all_labels) > 1:
                all_labels = [label for label in all_labels if label != '0']
                all_labels.sort()
            
            return ' '.join(all_labels)
        
    
    
def file_generator(data_folder, label_folder, label_mapper, parsers, root_folder=None):
    # determine root_folder if it was not already set (else it changes at each recursive function call) # TODO: move out of function
    root_folder = data_folder.split(os.sep)[-1] if not root_folder else root_folder
    
    # scan and iterate over the given folder
    with os.scandir(data_folder) as scan:
        for item in scan:
            filepath = item.path
            
            # if the item is a file
            if item.is_file():
                # extract file metadata
                file_info = get_metadata(filepath, root_folder)
                
                # parse and extract file content
                content = get_content(filepath, parsers)

                # yield file_info only if its content was parsed successfully
                if content:
                    file_info['content'] = content
    
                    # # determine file label
                    label = get_label(filepath, label_folder, label_mapper)
                    file_info['label'] = label
                    
                    yield file_info
                
            # else if the item is a folder, recursively call the function again on the filepath of this folder
            else:
                for subitem in file_generator(filepath, label_folder, label_mapper, parsers, root_folder):
                    yield subitem


def batch_generator(gen, batch_size):
    # initialize a list to contain the batch
    batch_list = []
    
    # perform generator 'next' call for 'batch_size' times
    for _ in range(batch_size):
        try: 
            print(_)
            batch_list.append(next(gen))
        except StopIteration:
            # generator 'next' call reaches end of iterator object
            yield batch_list
                
    yield batch_list


def folder_to_parquet(in_folder, label_folder, out_filename, batch_size=100, return_n_files=True):
    
    assert get_extension(out_filename) == '.parquet', "Output filename must have the extension: '.parquet'"
    print(f'Writing to file {out_filename} in batches of size {batch_size}...')
    
    parsers = {
        '.docx': doc_parser.load_docx, 
        '.pdf': doc_parser.load_pdf, 
        '.pptx': doc_parser.load_powerpoint, 
        'xlsx': doc_parser.load_excel,
        'xlsm': doc_parser.load_excel,
        '.msg': doc_parser.load_msg,
    }
    
    parquet_schema = pa.schema([
        ('filepath', pa.string()),
        ('filename', pa.string()),
        ('size', pa.string()),
        ('creation_date', pa.timestamp('s')),
        ('modified_date', pa.timestamp('s')),
        ('content', pa.string()),
        ('label', pa.string())
        # ('word_set', pa.string())
    ])

    pandas_schema = {'filepath': str,
                    'filename': str,
                    'size': str,
                    'creation_date': 'datetime64[s]',
                    'modified_date': 'datetime64[s]',
                    'content': str,
                    'label': str,
                    # 'word_set': str
                    }

    # label names are equal to the subfolder names in the label_folder
    label_classes = os.listdir(label_folder)
    assert all([os.path.isdir(os.path.join(label_folder, label)) for label in label_classes]), 'label_folder must contain only folders that specify the class label'

    # sort class names alphabetically and create a label_name --> label_id mapper
    label_classes.sort()
    label_mapper = {label: str(i+1) for i, label in enumerate(label_classes)}

    # initiaize file generator for the given folder
    file_gen = file_generator(in_folder, label_folder, label_mapper, parsers)
    
    total_rows = 0
    with pq.ParquetWriter(out_filename, parquet_schema, compression='gzip') as writer:
        
        while True:
            # process a batch of files
            batch = list(next(batch_generator(file_gen, batch_size=batch_size)))

            # if the batch contains anything
            if batch:
                # convert batch to pandas dataframe
                df = pd.DataFrame(batch)
                
                # convert batch dataframe to parquet table
                df = df.astype(pandas_schema)
                table = pa.Table.from_pandas(df)

                # write parquet table to parquet file
                writer.write_table(table)
                
                total_rows += table.num_rows
                print(f'Processed {total_rows} files...')

            # the batch doesn't contain anything, thus the function ends
            else:
                print(f'Successfully written {total_rows} files to {out_filename}!')
                
                if return_n_files:
                    return total_rows
                return

file_name = __file__.split('/')[-1]
path = __file__.replace(file_name, '')
print(path)
in_folder = path + "..\dataset\Lschijf_met_autorisatiematrix"
label_folder = path + "..\dataset\Lschijf_met_autorisatiematrix_gelabeld"
print(in_folder)
train_test_filename = 'train_test.parquet'
n_files = folder_to_parquet(in_folder, label_folder, train_test_filename, batch_size=500, return_n_files=True)

def batch_to_table(batch, pandas_schema):
    # convert batch to pandas dataframe
    df = batch.to_pandas()
    
    # # convert batch to parquet table
    df = df.astype(pandas_schema)
    table = pa.Table.from_pandas(df)
    return table


def train_test_split(parquet_file, n_files, train_size=0.7):
    parquet_schema = pa.schema([
        ('filepath', pa.string()),
        ('filename', pa.string()),
        ('size', pa.string()),
        ('creation_date', pa.timestamp('s')),
        ('modified_date', pa.timestamp('s')),
        ('content', pa.string()),
        ('label', pa.string())
        # ('word_set', pa.string())
    ])

    pandas_schema = {'filepath': str,
                     'filename': str,
                     'size': str,
                     'creation_date': 'datetime64[s]',
                     'modified_date': 'datetime64[s]',
                     'content': str,
                     'label': str,
                     # 'word_set': str
                    }
    
    train_total_rows = 0
    train_out_filename = 'train.parquet'
    with pq.ParquetWriter(train_out_filename, parquet_schema, compression='gzip') as writer:
        
        for batch in parquet_file.iter_batches(batch_size=1):
            
            # determine whether to start writing to test set
            if train_total_rows > int(n_files*train_size):
                break

            # convert batch to parquet tabel
            table = batch_to_table(batch, pandas_schema)
        
            # write parquet table to parquet file
            writer.write_table(table)
            
            train_total_rows += table.num_rows

    print(f'Successfully written {train_total_rows} training objects to {train_out_filename}!')

    test_total_rows = 0
    test_out_filename = 'test.parquet'
    with pq.ParquetWriter(test_out_filename, parquet_schema, compression='gzip') as writer:
        for i, batch in enumerate(parquet_file.iter_batches(batch_size=1)):
            if i < train_total_rows:
                continue

            # convert batch to parquet tabel
            table = batch_to_table(batch, pandas_schema)
        
            # write parquet table to parquet file
            writer.write_table(table)
            
            test_total_rows += table.num_rows

    print(f'Successfully written {test_total_rows} training objects to {test_out_filename}!')
    return train_out_filename, test_out_filename

train_test_parquet = pq.ParquetFile(train_test_filename)
train_filename, test_filename = train_test_split(train_test_parquet, n_files, train_size=0.7)

def parquet_generator(parquet_file, n_classes=4, batch_size=100, text_column='content', label_column='label'):
    # iterate over parquet file in batches
    for batch in parquet_file.iter_batches(batch_size=batch_size):
        # extract batch text and labels as a list
        batch_texts = batch[text_column].to_pylist()
        batch_labels = batch[label_column].to_pylist()   

        # remove multiple labels, always go for the label with the highest value
        # batch_labels = [int(item.split()[-1]) for item in batch_labels]

        
        # one-hot encode the labels
        batch_labels = [[1 if str(i) in item.split() else 0 for i in range(n_classes)] for item in batch_labels]

        # remove batch items that have no text (not needed if these are already discarded at parquet file generation)
        # batch_texts, batch_labels = list(zip(*[[text, label] for text, label in zip(batch_texts, batch_labels) if text]))
        # yield list(batch_texts), list(batch_labels)

        yield np.array(batch_texts), np.array(batch_labels)

def parquet_partial_fit(parquet_file, preprocessor, batch_size=100, train_column='content'):
    # initialize parquet_generator to iterate over parquet batches
    pg = parquet_generator(parquet_file, batch_size=batch_size, text_column='content', label_column='label')
    
    n_texts = 0
    # iterate over parquet batches
    for batch_texts, _ in pg:
        # perform a partial fit on the parquet batch
        preprocessor.partial_fit(batch_texts)

        n_texts += len(batch_texts)
        
    print(f'Preprocessor was fitted to {n_texts} texts...')    
    return preprocessor


# initialize preprocessor
preprocessor = HashingVectorizer(strip_accents='ascii')

# fit preprocessor on the parquet file, containing both the train and test set
preprocessor = parquet_partial_fit(train_test_parquet, preprocessor, batch_size=100)
preprocessor

class XGBoostIterator(xgboost.DataIter):
    def __init__(self, generator, preprocessor, target_class):
        self._generator = generator
        self._preprocessor = preprocessor
        self._it = 0
        self.target_class = target_class
        super().__init__()
    
    def next(self, input_data):
        try:
            texts, y = next(self._generator)
            
            X = self._preprocessor.transform(texts)
            
            # input_data(data=X, label=[y]) # for single label DMatrix generation

            # if self.target_class
            input_data(data=X, label=y[:,self.target_class]) # for one-hot encoded label DMatrix generation
            
            self._it += 1
            return 1
            
        except StopIteration:
            return 0 

    def reset(self):
        self._it = 0


def construct_DMatrix(parquet_filename, target_class):
    parquet_file = pq.ParquetFile(parquet_filename)
    
    # initialize parquet generator to iterate over train_parquet file in batches
    generator = parquet_generator(parquet_file, batch_size=1500)
    
    # initialize XGBoost iterator to train XGBoost in batches
    it = XGBoostIterator(generator, preprocessor, target_class)
    
    return xgboost.DMatrix(it)

def f1_eval(y_pred, data):
    y_true = data.get_label()
    f1 = f1_score(y_true, np.round(y_pred), average='macro')
    return 'f1', f1

from time import time
t1 = time()

params = {
    'objective': 'multi:softmax',
    'num_class': 2,
    # 'objective': 'binary:logistic',
    'tree_method': 'hist',
    'disable_default_eval_metric': True,
}

n_classes = 4
models = []
n_epochs = 10
for i in range(n_classes):
    print(f'Training model for class {i}...')
    Xy_train = construct_DMatrix(train_filename, target_class=i)
    Xy_test = construct_DMatrix(test_filename, target_class=i)

    best_model, best_score = None, 0
    for j in range(n_epochs):
    
        model = xgboost.train(params, Xy_train, 
                              custom_metric=f1_eval, 
                              num_boost_round=1,
                              early_stopping_rounds=1,
                              maximize=True,
                              xgb_model=(None if j == 0 else model),# Xgb model to be loaded before training (allows training continuation)
                              evals=[(Xy_train, 'train') , (Xy_test, 'valid')])
        score = model.best_score
        
        if score > best_score:
            best_model = model
            best_score = score
            
    models.append(best_model)
    print()

t2 = time()
print(f'All models were trained in {t2-t1:.0f}s')

f1_per_class = []
for i in range(n_classes):
    Xy_test = construct_DMatrix(test_filename, target_class=i)
    y_pred = models[i].predict(Xy_test)
    
    f1 = f1_eval(y_pred, Xy_test)[1]
    print(f'F1 for class {i}: {f1}')

    f1_per_class.append(f1)
    
print(f'\nMacro f1 = {np.mean(f1_per_class)}')

# dir(models[0])

# models[0].best_score
